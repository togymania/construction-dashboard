"""Monotekstroy — Yönetici Sunumu Oluşturucu

Kullanım:
    cd C:\\Projects\\construction-dashboard\\presentation
    python -m pip install python-pptx Pillow
    python build_presentation.py

Çıktı:
    Monotekstroy-Sunum.pptx (aynı klasörde)

Ekran görüntüleri:
    screenshots/ klasörüne aşağıdaki adlarla koy (varsa otomatik gömülür,
    yoksa slayt placeholder kutu gösterir):

      01-panel.png             - Panel (Ana sayfa, dashboard)
      02-projects.png          - Projeler listesi
      03-project-overview.png  - Proje detayı / Genel Bakış
      04-subcontractors.png    - Taşeronlar sekmesi
      05-workforce.png         - İşgücü sekmesi
      06-budget.png            - Bütçe sekmesi
      07-expenses.png          - Harcamalar sekmesi
      08-tenders-list.png      - İhaleler listesi
      09-tender-detail.png     - İhale detay (karşılaştırma grid)
      10-tender-ai.png         - AI Tender Analizi
      11-project-ai.png        - AI Proje Analizi
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Marka / tipografi
# ---------------------------------------------------------------------------

HERE = Path(__file__).parent
SHOTS = HERE / "screenshots"

# Indigo-Construct palette
PRIMARY = RGBColor(0x4F, 0x46, 0xE5)   # Indigo 600
SECONDARY = RGBColor(0x81, 0x8C, 0xF8) # Indigo 400
ACCENT = RGBColor(0xF5, 0x9E, 0x0B)    # Amber 500
DARK = RGBColor(0x0F, 0x17, 0x2A)      # Slate 900
SLATE800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE600 = RGBColor(0x47, 0x55, 0x69)
SLATE300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
ROSE = RGBColor(0xE1, 0x1D, 0x48)

FONT_HEAD = "Georgia"
FONT_BODY = "Calibri"

# 16:9 — 13.333 × 7.5 inç
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Yardımcı çizim fonksiyonları
# ---------------------------------------------------------------------------


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide, x, y, w, h, text, *,
    size=14, bold=False, color=SLATE800,
    font=FONT_BODY, align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP, italic=False,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, bullets, *,
                size=14, color=SLATE800, font=FONT_BODY,
                bullet_color=PRIMARY, spacing=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        # bullet glyph
        r_bullet = p.add_run()
        r_bullet.text = "▸ "
        r_bullet.font.name = font
        r_bullet.font.size = Pt(size)
        r_bullet.font.bold = True
        r_bullet.font.color.rgb = bullet_color
        # text
        r = p.add_run()
        r.text = b
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return tb


def add_image_or_placeholder(slide, x, y, w, h, filename, caption=None):
    """Eğer ekran görüntüsü dosyası varsa göm, yoksa gri placeholder kutu."""
    path = SHOTS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        # Placeholder
        box = add_rect(slide, x, y, w, h, SLATE100, line=SLATE300)
        add_text(
            slide,
            x + Inches(0.2), y + h / 2 - Inches(0.3),
            w - Inches(0.4), Inches(0.6),
            f"[ Ekran görüntüsü: {filename} ]",
            size=12, italic=True, color=SLATE600,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
    if caption:
        add_text(
            slide,
            x, y + h + Inches(0.05),
            w, Inches(0.3),
            caption,
            size=10, italic=True, color=SLATE600,
            align=PP_ALIGN.CENTER,
        )


def add_pageno(slide, n, total):
    add_text(
        slide,
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.45),
        Inches(1.0), Inches(0.3),
        f"{n} / {total}",
        size=9, color=SLATE600,
        align=PP_ALIGN.RIGHT,
    )


def add_brand_strip(slide):
    """Sol kenarda 6pt dikey indigo şerit — visual motif."""
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, PRIMARY)


def add_title(slide, kicker, title, color=DARK):
    if kicker:
        add_text(
            slide,
            Inches(0.5), Inches(0.45),
            Inches(8), Inches(0.3),
            kicker.upper(),
            size=11, bold=True, color=PRIMARY,
            font=FONT_BODY,
        )
    add_text(
        slide,
        Inches(0.5), Inches(0.75),
        Inches(12), Inches(0.9),
        title,
        size=32, bold=True, color=color,
        font=FONT_HEAD,
    )


# ---------------------------------------------------------------------------
# Slayt builders
# ---------------------------------------------------------------------------


def slide_cover(prs, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Koyu zemin
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    # Sağ tarafta amber accent geometri
    add_rect(s, SLIDE_W - Inches(2.8), Inches(0), Inches(2.8), SLIDE_H, PRIMARY)
    add_rect(s, SLIDE_W - Inches(0.5), Inches(0), Inches(0.5), SLIDE_H, ACCENT)

    add_text(
        s, Inches(0.7), Inches(1.0), Inches(9), Inches(0.5),
        "MONOTEKSTROY",
        size=14, bold=True, color=ACCENT, font=FONT_BODY,
    )
    add_text(
        s, Inches(0.7), Inches(1.6), Inches(9), Inches(2.2),
        "AI Destekli\nİnşaat Proje\nYönetim Platformu",
        size=54, bold=True, color=WHITE, font=FONT_HEAD,
    )
    add_text(
        s, Inches(0.7), Inches(4.7), Inches(9), Inches(1.0),
        "Tender analizinden bütçe takibine, taşeron yönetiminden\nişgücü analitiğine kadar tüm proje süreçleri tek dashboard'da",
        size=18, color=SLATE300, font=FONT_BODY, italic=True,
    )
    # Alt blok
    add_rect(s, Inches(0.7), Inches(6.4), Inches(0.4), Inches(0.4), ACCENT)
    add_text(
        s, Inches(1.2), Inches(6.4), Inches(8), Inches(0.4),
        "Yönetici Sunumu  ·  Mayıs 2026",
        size=14, color=SLATE300, font=FONT_BODY,
    )


def slide_vision(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "Neden Monotekstroy?", "Yönetici Karar Merkezi")

    add_text(
        s, Inches(0.5), Inches(1.85), Inches(7.5), Inches(1.6),
        "İnşaat şirketinin günlük operasyonel sorularını —\n"
        "“Hangi taşeron geç kalıyor?”, “Bütçe nerede aşılıyor?”,\n"
        "“Hangi teklif daha iyi?” — tek bir dashboard'da, AI destekli\n"
        "analizle yanıtlayan bütünleşik bir platform.",
        size=16, color=SLATE800,
    )

    # 4 ana sütun KPI kartları
    cards = [
        ("Tender Süresi",      "2 saat → 5 dk",   "PDF/Excel teklif → AI ile karşılaştırma tablosu", PRIMARY),
        ("Otomatik Sınıflama", "%70–80",          "Banka ekstresi satırlarının AI ile bütçe koduna eşlenmesi", GREEN),
        ("Erken Uyarı",        "EAC + CPI",       "Bütçe aşımı geri dönüşsüz olmadan önce tespit", ACCENT),
        ("Pazarlık Geçmişi",   "v1 → v2 → v3",    "Aynı tedarikçinin her revizesi otomatik tarihçeye işlenir", SECONDARY),
    ]
    card_y = Inches(4.0)
    card_h = Inches(2.6)
    card_w = Inches(3.0)
    gap = Inches(0.2)
    total_w = card_w * 4 + gap * 3
    start_x = (SLIDE_W - total_w) / 2
    for i, (label, big, desc, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, card_y, card_w, card_h, WHITE, line=SLATE300)
        add_rect(s, x, card_y, Inches(0.12), card_h, color)
        add_text(s, x + Inches(0.3), card_y + Inches(0.25), card_w, Inches(0.3),
                 label.upper(), size=10, bold=True, color=SLATE600, font=FONT_BODY)
        add_text(s, x + Inches(0.3), card_y + Inches(0.6), card_w, Inches(0.9),
                 big, size=28, bold=True, color=color, font=FONT_HEAD)
        add_text(s, x + Inches(0.3), card_y + Inches(1.55), card_w - Inches(0.5), Inches(1.0),
                 desc, size=12, color=SLATE800)

    add_pageno(s, n, total)


def slide_architecture(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "Teknik Mimari", "Modern Bulut Yığını + Anthropic Claude")

    # Sol: stack tablosu
    rows = [
        ("Frontend",  "Next.js 16 · React 19 · Tailwind",  "Vercel (Auto-Deploy)"),
        ("Backend",   "FastAPI 0.115 · Python 3.12",       "Render (Frankfurt)"),
        ("Veritabanı","PostgreSQL 16",                      "Neon (Frankfurt)"),
        ("AI Motoru", "Anthropic Claude Sonnet 4.5",        "İki dilli (RU / TR)"),
        ("Dosya İşleme","pdfplumber · openpyxl",            "PDF + Excel → JSON"),
        ("Güvenlik",  "JWT auth · CORS · RBAC",             "Owner / PM / Viewer"),
    ]
    table_x = Inches(0.5)
    table_y = Inches(1.85)
    row_h = Inches(0.55)
    col_w = [Inches(2.2), Inches(3.8), Inches(2.5)]
    headers = ["Katman", "Teknoloji", "Çalıştığı Yer"]
    # header
    hx = table_x
    for i, ht in enumerate(headers):
        add_rect(s, hx, table_y, col_w[i], row_h, PRIMARY)
        add_text(s, hx + Inches(0.15), table_y, col_w[i], row_h,
                 ht, size=12, bold=True, color=WHITE, font=FONT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)
        hx += col_w[i]
    # rows
    for ri, row in enumerate(rows):
        ry = table_y + row_h * (ri + 1)
        bg = WHITE if ri % 2 == 0 else SLATE100
        rx = table_x
        for ci, cell in enumerate(row):
            add_rect(s, rx, ry, col_w[ci], row_h, bg, line=SLATE300)
            add_text(s, rx + Inches(0.15), ry, col_w[ci], row_h,
                     cell, size=11, color=SLATE800,
                     bold=(ci == 0), anchor=MSO_ANCHOR.MIDDLE)
            rx += col_w[ci]

    # Sağ: AI vurgu kutusu
    box_x = Inches(9.3)
    box_y = Inches(1.85)
    box_w = Inches(3.5)
    box_h = Inches(5.2)
    add_rect(s, box_x, box_y, box_w, box_h, DARK)
    add_rect(s, box_x, box_y, box_w, Inches(0.12), ACCENT)
    add_text(s, box_x + Inches(0.3), box_y + Inches(0.4), box_w, Inches(0.4),
             "AI YETENEKLERİ", size=11, bold=True, color=ACCENT, font=FONT_BODY)
    add_text(s, box_x + Inches(0.3), box_y + Inches(0.85), box_w - Inches(0.5), Inches(0.7),
             "Claude Sonnet 4.5", size=22, bold=True, color=WHITE, font=FONT_HEAD)
    ai_features = [
        "Yapılandırılmış JSON çıkarımı",
        "PDF / Excel okuma",
        "RU / TR çift dilli üretim",
        "Risk analizi ve öneri",
        "Otomatik sınıflandırma",
        "Cascading fallback motoru",
    ]
    add_bullets(s, box_x + Inches(0.3), box_y + Inches(1.9),
                box_w - Inches(0.5), Inches(3.0),
                ai_features, size=12, color=SLATE100, bullet_color=ACCENT)

    add_pageno(s, n, total)


def slide_tab(prs, n, total, kicker, title, screenshot, bullets, summary):
    """Ortak sekme slaydı: sol açıklama + sağ ekran görüntüsü."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, kicker, title)

    # Sol: özet + bullet'lar
    add_text(
        s, Inches(0.5), Inches(1.85),
        Inches(5.5), Inches(1.6),
        summary, size=14, color=SLATE800,
    )
    add_bullets(
        s, Inches(0.5), Inches(3.7),
        Inches(5.5), Inches(3.4),
        bullets, size=14, color=SLATE800, spacing=8,
    )

    # Sağ: ekran görüntüsü
    img_x = Inches(6.4)
    img_y = Inches(1.85)
    img_w = Inches(6.5)
    img_h = Inches(4.5)
    add_image_or_placeholder(s, img_x, img_y, img_w, img_h, screenshot)

    add_pageno(s, n, total)


def slide_tender_showcase(prs, n, total):
    """İhale Detayı — sunumun yıldız ekranı, biraz farklı layout."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "İhaleler — Detay Ekranı", "Sunumun Yıldız Ekranı")

    # Üst — tek satır info kartları
    kpis = [
        ("Hiyerarşik satırlar", "1, 1.1, 1.2", PRIMARY),
        ("Varyantlar", "Aynı firma, farklı ürün", SECONDARY),
        ("Revizyon takibi", "v1 → v2 → v3", ACCENT),
        ("Renk gradyanı", "Yeşil = En ucuz", GREEN),
    ]
    kx = Inches(0.5)
    ky = Inches(1.85)
    kw = Inches(3.0)
    kh = Inches(0.9)
    gap = Inches(0.12)
    for i, (label, value, col) in enumerate(kpis):
        x = kx + i * (kw + gap)
        add_rect(s, x, ky, kw, kh, WHITE, line=SLATE300)
        add_rect(s, x, ky, Inches(0.1), kh, col)
        add_text(s, x + Inches(0.25), ky + Inches(0.12), kw, Inches(0.25),
                 label.upper(), size=9, bold=True, color=SLATE600, font=FONT_BODY)
        add_text(s, x + Inches(0.25), ky + Inches(0.4), kw, Inches(0.45),
                 value, size=14, bold=True, color=SLATE800, font=FONT_HEAD)

    # Büyük ekran görüntüsü
    add_image_or_placeholder(
        s, Inches(0.5), Inches(3.0), Inches(12.3), Inches(3.7),
        "09-tender-detail.png",
        caption="Резиновое напольное покрытие · 2 firma · ООО АгроЦентрик'in iki varyantı (Dairy Plus / Terras)",
    )

    add_pageno(s, n, total)


def slide_ai_tender(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "AI Modülü 1/2", "İhale Analizi — 6 Bölümlü Karar Raporu")

    # Sol: 6 bölüm tile grid
    sections = [
        ("Genel Bakış",   "Bid sayısı, ortalama, en düşük/en yüksek, dağılım yüzdesi"),
        ("Karşılaştırma", "Her firma için fiyat ve teslim süresi"),
        ("Analiz",        "En ucuz, en hızlı, en dengeli firma"),
        ("Riskler",       "Her firma için olası problemler"),
        ("Tavsiye",       "Seçilen firma + alternatif + güven yüzdesi"),
        ("Özet Anlatım",  "Türkçe paragraf — yöneticiye 3 satırda durum"),
    ]
    sx = Inches(0.5)
    sy = Inches(1.85)
    sw = Inches(5.7)
    sh_h = Inches(0.7)
    for i, (title, desc) in enumerate(sections):
        y = sy + i * (sh_h + Inches(0.12))
        add_rect(s, sx, y, sw, sh_h, SLATE100, line=SLATE300)
        # numara dairesi
        add_rect(s, sx + Inches(0.15), y + Inches(0.15), Inches(0.4), Inches(0.4), PRIMARY)
        add_text(s, sx + Inches(0.15), y + Inches(0.15),
                 Inches(0.4), Inches(0.4),
                 str(i + 1), size=14, bold=True, color=WHITE,
                 font=FONT_HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, sx + Inches(0.7), y + Inches(0.08), sw, Inches(0.3),
                 title, size=13, bold=True, color=SLATE800, font=FONT_HEAD)
        add_text(s, sx + Inches(0.7), y + Inches(0.35), sw - Inches(0.7), Inches(0.35),
                 desc, size=11, color=SLATE600)

    # Sağ: ekran görüntüsü
    add_image_or_placeholder(
        s, Inches(6.6), Inches(1.85),
        Inches(6.3), Inches(4.7),
        "10-tender-ai.png",
        caption="Ластик Каплама teklifi için AI önerisi: %75 güvenle Дэйри Плюс",
    )

    add_pageno(s, n, total)


def slide_ai_project(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "AI Modülü 2/2", "Proje Geneli Yönetici Özeti")

    add_text(
        s, Inches(0.5), Inches(1.85),
        Inches(12.3), Inches(0.8),
        "Tek tıkla 5 modülü sentezleyen executive dashboard — sahada ne oluyor, ne yanlış,"
        " ne acil. Her modül Türkçe AI yorumuyla birlikte gelir.",
        size=14, color=SLATE800, italic=True,
    )

    modules = [
        ("Subcontractor & Schedule", "Geciken sözleşmeler, kritik firmalar, disiplin kırılımı", ROSE),
        ("Data Quality",             "Sınıflandırılmamış kayıtlar, atanmamış ödemeler, kirli oran",  ACCENT),
        ("Financial (EAC)",           "BAC, AC, EV, CPI, EAC, VAC — bütçe sapma tahmini",            GREEN),
        ("Workforce & Productivity",  "Headcount, adam-saat, verimlilik, sapma",                     SECONDARY),
        ("Risk Analysis",             "Üç temel risk + her birinin etki + nedeni",                   PRIMARY),
    ]
    mx = Inches(0.5)
    my = Inches(2.85)
    mw = Inches(2.4)
    mh = Inches(2.5)
    gap = Inches(0.15)
    for i, (title, desc, color) in enumerate(modules):
        x = mx + i * (mw + gap)
        add_rect(s, x, my, mw, mh, WHITE, line=SLATE300)
        add_rect(s, x, my, mw, Inches(0.15), color)
        add_text(s, x + Inches(0.2), my + Inches(0.35),
                 mw - Inches(0.4), Inches(0.7),
                 title, size=12, bold=True, color=SLATE800, font=FONT_HEAD)
        add_text(s, x + Inches(0.2), my + Inches(1.15),
                 mw - Inches(0.4), Inches(1.3),
                 desc, size=10, color=SLATE600)

    # Alt: ekran görüntüsü
    add_image_or_placeholder(
        s, Inches(0.5), Inches(5.7),
        Inches(12.3), Inches(1.5),
        "11-project-ai.png",
        caption="AI Proje Analizi · Türkçe çıktı · Refresh ile yenilenir",
    )

    add_pageno(s, n, total)


def slide_impact(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # Koyu zemin
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(s, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT)

    add_text(s, Inches(0.5), Inches(0.5), Inches(8), Inches(0.4),
             "BEKLENEN ETKİ", size=12, bold=True, color=ACCENT, font=FONT_BODY)
    add_text(s, Inches(0.5), Inches(0.95), Inches(12), Inches(0.9),
             "Operasyonel Kazanımlar", size=32, bold=True, color=WHITE, font=FONT_HEAD)

    impacts = [
        ("2 saat → 5 dakika",     "Bir tender'ın PDF'ten karşılaştırma tablosuna geçişi"),
        ("%70 – 80",              "Banka ekstresi sınıflandırmasında otomasyon oranı"),
        ("Sıfır manuel transfer", "PDF/Excel'den AI doğrudan veritabanına yazıyor"),
        ("Erken uyarı",           "EAC + CPI bütçe aşımını geri dönüşsüz olmadan tespit"),
        ("Tek doğruluk kaynağı",  "Tüm portföy — projeler, taşeronlar, işgücü — tek panelde"),
        ("Pazarlık şeffaflığı",   "Her tedarikçi revizesi otomatik history'e işleniyor"),
    ]
    cy = Inches(2.3)
    ch = Inches(0.78)
    gap = Inches(0.12)
    for i, (big, desc) in enumerate(impacts):
        y = cy + i * (ch + gap)
        # accent left line
        add_rect(s, Inches(0.5), y, Inches(0.08), ch, ACCENT)
        add_text(s, Inches(0.8), y, Inches(4.0), ch,
                 big, size=20, bold=True, color=WHITE, font=FONT_HEAD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(5.0), y, Inches(8.0), ch,
                 desc, size=14, color=SLATE300,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.5), SLIDE_H - Inches(0.45), Inches(8), Inches(0.3),
             f"{n} / {total}", size=9, color=SLATE300,
             align=PP_ALIGN.LEFT)


def slide_roadmap(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_brand_strip(s)
    add_title(s, "Yol Haritası", "Sunum Sonrası Genişleme")

    phases = [
        ("Şu An", "Canlıda",
         ["6 modül + AI", "Üretim ortamında çalışıyor", "Demo verisi yüklü"],
         GREEN, "✓"),
        ("Faz 1", "1-2 hafta",
         ["TDF formatında Excel export", "Schedule (Gantt-lite)", "Risk Register modülü"],
         PRIMARY, "→"),
        ("Faz 2", "3-4 hafta",
         ["Market price web search", "Mobile responsive iyileştirme", "PDF export raporları"],
         SECONDARY, "→"),
        ("Faz 3", "Vizyon",
         ["Real-time WebSocket güncellemeler", "Multi-language genişletme", "Yapay zekâ ile sözleşme inceleme"],
         ACCENT, "★"),
    ]
    px = Inches(0.5)
    py = Inches(1.85)
    pw = Inches(3.05)
    ph = Inches(4.7)
    gap = Inches(0.18)
    for i, (title, time, items, col, glyph) in enumerate(phases):
        x = px + i * (pw + gap)
        add_rect(s, x, py, pw, ph, WHITE, line=SLATE300)
        add_rect(s, x, py, pw, Inches(0.8), col)
        # glyph
        add_text(s, x + Inches(0.25), py + Inches(0.1),
                 Inches(0.5), Inches(0.5),
                 glyph, size=22, bold=True, color=WHITE, font=FONT_HEAD)
        add_text(s, x + Inches(0.85), py + Inches(0.12),
                 pw, Inches(0.3),
                 title.upper(), size=11, bold=True, color=WHITE, font=FONT_BODY)
        add_text(s, x + Inches(0.85), py + Inches(0.4),
                 pw, Inches(0.35),
                 time, size=14, bold=True, color=WHITE, font=FONT_HEAD)
        # body
        add_bullets(s, x + Inches(0.25), py + Inches(1.05),
                    pw - Inches(0.5), Inches(3.5),
                    items, size=12, color=SLATE800,
                    bullet_color=col, spacing=8)

    add_pageno(s, n, total)


def slide_thanks(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(s, SLIDE_W - Inches(3.5), Inches(0), Inches(3.5), SLIDE_H, PRIMARY)
    add_rect(s, SLIDE_W - Inches(0.5), Inches(0), Inches(0.5), SLIDE_H, ACCENT)

    add_text(s, Inches(0.7), Inches(1.2), Inches(8), Inches(0.5),
             "SUNUM SONU", size=12, bold=True, color=ACCENT, font=FONT_BODY)
    add_text(s, Inches(0.7), Inches(1.8), Inches(9), Inches(1.5),
             "Teşekkürler", size=72, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(3.5), Inches(8), Inches(0.6),
             "Sorularınız?", size=24, color=SLATE300, font=FONT_HEAD, italic=True)

    add_text(s, Inches(0.7), Inches(5.5), Inches(9), Inches(0.4),
             "Topal Tolga", size=18, bold=True, color=WHITE, font=FONT_HEAD)
    add_text(s, Inches(0.7), Inches(6.0), Inches(9), Inches(0.4),
             "tolgatopal1999@gmail.com", size=13, color=SLATE300, font=FONT_BODY)
    add_text(s, Inches(0.7), Inches(6.4), Inches(9), Inches(0.4),
             "Canlı Demo: monotek-stroy-pm.vercel.app", size=13, color=SLATE300, font=FONT_BODY)


# ---------------------------------------------------------------------------
# Ana akış
# ---------------------------------------------------------------------------


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 16  # toplam slayt

    slide_cover(prs, total)
    slide_vision(prs, 2, total)
    slide_architecture(prs, 3, total)

    # 4. Panel
    slide_tab(
        prs, 4, total,
        "Modül 1 — Panel",
        "Yönetici Paneli (Ana Sayfa)",
        "01-panel.png",
        [
            "4 KPI özeti: aktif proje, toplam bütçe, on-track, açık riskler",
            "Data Quality kartı — kirli kayıt sayısı, hemen temizleme linki",
            "Son aktiviteler akışı — gerçek zamanlı proje hareketleri",
            "Aktif projeler listesi — durumlarıyla birlikte",
            "Tek bakışta tüm portföyün sağlığı",
        ],
        "Yöneticinin sabah günü başlatırken ilk göreceği ekran. "
        "Anlık bir nabız — neyin acil olduğunu 5 saniyede gösterir.",
    )
    # 5. Projeler
    slide_tab(
        prs, 5, total,
        "Modül 2 — Projeler",
        "Portföy Görünümü",
        "02-projects.png",
        [
            "Tüm projelerin listesi: ad, durum, sağlık, lokasyon",
            "Bütçe / kullanım / ilerleme — tek tablo, tek bakış",
            "Filtreleme ve arama (firma, lokasyon, durum)",
            "Yeni Proje sihirbazıyla hızlı kayıt",
            "Her satır tıklanabilir — detay sayfasına geçiş",
        ],
        "Aktif portföyün özet listesi. CEO'nun “bugün hangi projeler ileri, hangileri geri?” "
        "sorusunun anında yanıtı.",
    )
    # 6. Proje detayı
    slide_tab(
        prs, 6, total,
        "Modül 3 — Proje Detayı",
        "Genel Bakış + EAC Forecast",
        "03-project-overview.png",
        [
            "Header: durum + sağlık + zaman çizelgesi + bütçe",
            "EAC Forecast: BAC, AC, EAC, Variance — Earned Value Management",
            "Modül kartları: Taşeron, İşgücü, Bütçe, Harcamalar, Takvim, Risk, Rapor, İhale, AI",
            "Sol kenar navigasyonu — projenin her detayına 1 tıkla",
            "Tek sayfada 360° proje özeti",
        ],
        "Her proje için ayrı bir “mini-platform”. Soldaki menüden 9 modüle anında geçiş, "
        "üstte EAC forecast ile bütçe sapması canlı.",
    )
    # 7. Taşeronlar
    slide_tab(
        prs, 7, total,
        "Modül 4 — Taşeronlar",
        "Sözleşme, Ödeme ve Nakit Akışı Tahmini",
        "04-subcontractors.png",
        [
            "5 KPI: toplam taşeron, aktif sözleşme, geciken, toplam değer, ödeme ilerlemesi",
            "Ödeme durum dağılımı (ödendi / beklemede / geciken)",
            "Değere göre en büyük 5 taşeron — rank grafiği",
            "Aylık ödemeler trendi (son 6 ay)",
            "AI Nakit Akışı — 3 aylık tahmin (En iyi / Gerçekleşen / Olası / En kötü)",
        ],
        "Taşeronlar ile olan tüm finansal ilişkinin tek panelde toplu yönetimi. "
        "AI nakit akışı tahmini ile ileriye dönük risk planlaması.",
    )
    # 8. İşgücü
    slide_tab(
        prs, 8, total,
        "Modül 5 — İşgücü",
        "Sahadaki Personel Analitiği",
        "05-workforce.png",
        [
            "4 KPI: toplam, direkt, taşeron, haftalık değişim yüzdesi",
            "Today by Company — her firmanın direkt/dolaylı/taşeron kırılımı",
            "Günlük işgücü trendi grafiği (son 9 gün)",
            "AI Insights: günlük yorumlar — “Bugün toplam 115 azaldı”, “Civil 1280 işçi ile %74”",
            "Excel ile çoklu firma veri import",
        ],
        "Sahadaki gerçek insan sayısı, disiplin kırılımı, trend — AI'nin günlük yorumuyla. "
        "“Bugün kaç kişi sahada?” sorusunun bilimsel cevabı.",
    )
    # 9. Bütçe
    slide_tab(
        prs, 9, total,
        "Modül 6 — Bütçe",
        "Planlama vs Gerçekleşme",
        "06-budget.png",
        [
            "4 KPI: proje bütçesi, planlanan, harcanan, kullanım yüzdesi",
            "Planned by Category — pasta grafik (Bina, Yollar, Elektrik...)",
            "Planned vs Spent — kategori bazlı bar grafik karşılaştırması",
            "Budget Items / Planned vs Actual / Subcontractors — 3 sekmeli detay",
            "Excel ile bütçe kalemi import",
        ],
        "Onaylı bütçe → kategori kırılımı → harcama gerçekleşmesi tek görselde. "
        "16 milyar ruble bütçeli projenin nereye gittiğini takip eder.",
    )
    # 10. Harcamalar
    slide_tab(
        prs, 10, total,
        "Modül 7 — Harcamalar",
        "Banka Ekstresi → AI ile Otomatik Sınıflandırma",
        "07-expenses.png",
        [
            "3 KPI: toplam gelir, gider, net",
            "9.258 satır → bütçe kodlarına otomatik eşleme",
            "Status etiketleri: bütçe kodsuz / taşeronsuz uyarıları",
            "Filtre: Tümü / Gelir / Gider / Bütçe Kodsuz / Taşeronsuz",
            "Excel/CSV import → AI otomatik sınıflandırma %70-80 başarı",
        ],
        "Banka ekstresini doğrudan içe aktarın — AI her satırı bütçe kodu ile firma "
        "eşleştirip işaretler. Manuel veri girişini neredeyse sıfıra indirir.",
    )

    # 11. İhale showcase
    slide_tender_showcase(prs, 11, total)

    # 12. AI Tender Analizi
    slide_ai_tender(prs, 12, total)

    # 13. AI Proje Analizi
    slide_ai_project(prs, 13, total)

    # 14. Etki
    slide_impact(prs, 14, total)
    # 15. Yol haritası
    slide_roadmap(prs, 15, total)
    # 16. Teşekkürler
    slide_thanks(prs, 16, total)

    out = HERE / "Monotekstroy-Sunum.pptx"
    prs.save(str(out))
    print(f"✓ Sunum oluşturuldu: {out}")
    print(f"  Toplam {len(prs.slides)} slayt")
    if SHOTS.exists():
        present = sorted([p.name for p in SHOTS.glob("*.png")])
        if present:
            print(f"  Eklenmiş ekran görüntüleri: {len(present)} adet")
        else:
            print("  (screenshots/ klasörü boş — placeholder kutuları kullanıldı)")
    else:
        print("  (screenshots/ klasörü yok — placeholder kutuları kullanıldı)")
        print("  Ekran görüntülerini eklemek için README.md'ye bak.")


if __name__ == "__main__":
    build()
